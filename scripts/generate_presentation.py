#!/usr/bin/env python3
"""Generate USPTO demo PowerPoint deck for AWS Health Notifications Tracker."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR

# ── Brand colours ──────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x23, 0x2F, 0x3E)
ORANGE = RGBColor(0xFF, 0x99, 0x00)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF5, 0xF5, 0xF5)
DGRAY  = RGBColor(0x44, 0x44, 0x44)
RED    = RGBColor(0xD1, 0x3F, 0x28)
GREEN  = RGBColor(0x1E, 0x8E, 0x4C)
BLUE   = RGBColor(0x00, 0x73, 0xBB)

SCREENS = os.path.join(os.path.dirname(__file__), "..", "mock_screens")
OUT     = os.path.join(os.path.dirname(__file__), "..", "presentation", "aws-health-tracker-deck.pptx")
os.makedirs(os.path.dirname(OUT), exist_ok=True)

W, H = Inches(13.33), Inches(7.5)   # 16:9 widescreen

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank


# ── helpers ────────────────────────────────────────────────────────────────────
def add_rect(slide, left, top, width, height, fill=None, line=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE=1
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, size=18,
                bold=False, color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return tb

def add_bullet_box(slide, left, top, width, height, bullets, size=16,
                   color=DGRAY, bold_first=False, spacing=1.2):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.space_after  = Pt(4)
        run = p.add_run()
        run.text = b
        run.font.size  = Pt(size)
        run.font.color.rgb = color
        run.font.bold  = (bold_first and i == 0)
    return tb

def navy_header(slide, title, subtitle=None):
    """Full-width navy header bar at top."""
    add_rect(slide, 0, 0, W, Inches(1.2), fill=NAVY)
    add_textbox(slide, Inches(0.4), Inches(0.18), Inches(10), Inches(0.65),
                title, size=28, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, Inches(0.4), Inches(0.75), Inches(10), Inches(0.4),
                    subtitle, size=14, color=ORANGE)
    # orange accent line
    add_rect(slide, 0, Inches(1.2), W, Inches(0.06), fill=ORANGE)

def footer(slide):
    add_rect(slide, 0, Inches(7.1), W, Inches(0.4), fill=NAVY)
    add_textbox(slide, Inches(0.3), Inches(7.12), Inches(6), Inches(0.3),
                "AWS Health Notifications Tracker  |  USPTO AWS Innovation Team",
                size=9, color=RGBColor(0xAA,0xAA,0xAA))
    add_textbox(slide, Inches(10.5), Inches(7.12), Inches(2.5), Inches(0.3),
                "CONFIDENTIAL", size=9, color=RGBColor(0xAA,0xAA,0xAA), align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=NAVY)
# orange side accent
add_rect(s, 0, 0, Inches(0.18), H, fill=ORANGE)
# large title
add_textbox(s, Inches(0.6), Inches(1.6), Inches(11), Inches(1.4),
            "AWS Health Notifications Tracker",
            size=42, bold=True, color=WHITE)
# orange rule
add_rect(s, Inches(0.6), Inches(3.15), Inches(9), Inches(0.07), fill=ORANGE)
# subtitle
add_textbox(s, Inches(0.6), Inches(3.35), Inches(10), Inches(0.6),
            "Centralized Visibility & Actionable Intelligence for 200+ Account AWS Organizations",
            size=20, color=RGBColor(0xCC,0xCC,0xCC))
# presenter block
add_textbox(s, Inches(0.6), Inches(4.6), Inches(8), Inches(0.4),
            "Demo  |  USPTO AWS Innovation Team  |  April 2026",
            size=14, color=ORANGE)
# AWS badge
badge = add_rect(s, Inches(10.6), Inches(6.5), Inches(2.3), Inches(0.7), fill=ORANGE)
add_textbox(s, Inches(10.65), Inches(6.52), Inches(2.2), Inches(0.65),
            "Powered by AWS", size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
navy_header(s, "The Challenge", "Managing AWS Health at scale across a large AWS Organization")
footer(s)

problems = [
    ("⚠  Missed Deadlines",
     "Critical maintenance windows and deprecation deadlines buried in individual account dashboards — "
     "no org-wide view means teams miss required actions."),
    ("📋  Manual, Fragmented Tracking",
     "Operations teams manually check 200+ accounts in the AWS Console. "
     "Spreadsheets used to track follow-up status become stale within hours."),
    ("🔇  No Prioritization or Intelligence",
     "Raw AWS Health events lack context. Teams cannot distinguish Critical from Low urgency "
     "without reading every event in full — and there is no owner assignment workflow."),
]

y = Inches(1.45)
for icon_title, body in problems:
    add_rect(s, Inches(0.4), y, Inches(12.4), Inches(1.35), fill=LGRAY, line=RGBColor(0xDD,0xDD,0xDD))
    add_textbox(s, Inches(0.6), y + Inches(0.12), Inches(11.5), Inches(0.4),
                icon_title, size=16, bold=True, color=NAVY)
    add_textbox(s, Inches(0.6), y + Inches(0.5), Inches(11.5), Inches(0.7),
                body, size=13, color=DGRAY, wrap=True)
    y += Inches(1.5)

add_textbox(s, Inches(0.4), Inches(6.55), Inches(12), Inches(0.4),
            "Result: Compliance risk, reactive operations, and engineer burnout from manual monitoring.",
            size=13, bold=True, color=RED, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Solution Overview
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
navy_header(s, "Solution: AWS Health Notifications Tracker",
            "Automated aggregation · AI summarization · Deadline tracking · Follow-up workflow")
footer(s)

cols = [
    ("🔄  Automated Collection",   "EventBridge triggers Health Collector Lambda every 15 min across all 200 accounts via AWS Organizations Health API."),
    ("🤖  AI Summarization",       "Amazon Bedrock (Claude Haiku) converts raw events into plain-English summaries, action items, and urgency scores."),
    ("📊  Centralized Dashboard",  "React SPA served via CloudFront shows stats, sortable/filterable event list, and deadline countdowns."),
    ("🔔  Proactive Reminders",    "Daily SNS digest emails for deadlines within 7 days — Critical events escalate at 3 days."),
    ("✅  Follow-up Workflow",     "Per-event owner assignment, status tracking (Pending → In Progress → Resolved), and notes — all persisted in DynamoDB."),
]

x = Inches(0.25)
for title, body in cols:
    add_rect(s, x, Inches(1.45), Inches(2.45), Inches(4.8), fill=WHITE, line=RGBColor(0xDD,0xDD,0xDD))
    add_rect(s, x, Inches(1.45), Inches(2.45), Inches(0.45), fill=NAVY)
    add_textbox(s, x + Inches(0.1), Inches(1.48), Inches(2.25), Inches(0.4),
                title[:2], size=18, color=ORANGE)
    add_textbox(s, x + Inches(0.1), Inches(2.0), Inches(2.25), Inches(0.45),
                title[3:], size=13, bold=True, color=NAVY)
    add_textbox(s, x + Inches(0.1), Inches(2.52), Inches(2.25), Inches(3.5),
                body, size=11, color=DGRAY, wrap=True)
    x += Inches(2.6)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Architecture Diagram
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
navy_header(s, "Architecture", "Serverless · FedRAMP Authorized · us-east-1 · ~$10–15/month")
footer(s)

# ── draw boxes + arrows ────────────────────────────────────────────────────────
def box(sl, left, top, w, h, label, sub="", bg=NAVY, fg=WHITE, sub_color=ORANGE):
    add_rect(sl, left, top, w, h, fill=bg)
    add_textbox(sl, left+Inches(0.06), top+Inches(0.06), w-Inches(0.12), Inches(0.35),
                label, size=11, bold=True, color=fg, align=PP_ALIGN.CENTER)
    if sub:
        add_textbox(sl, left+Inches(0.06), top+Inches(0.38), w-Inches(0.12), Inches(0.28),
                    sub, size=9, color=sub_color, align=PP_ALIGN.CENTER)

def arrow(sl, x1, y1, x2, y2):
    from pptx.util import Pt
    connector = sl.shapes.add_connector(1, x1, y1, x2, y2)
    connector.line.color.rgb = ORANGE
    connector.line.width = Pt(1.8)

BW, BH = Inches(1.85), Inches(0.72)

# Row 1 — triggers
box(s, Inches(0.3),  Inches(1.5),  BW, BH, "EventBridge", "rate(15 min)",    bg=RGBColor(0x7B,0x2D,0x8B))
box(s, Inches(0.3),  Inches(3.0),  BW, BH, "EventBridge", "cron(0 9 * * ?)", bg=RGBColor(0x7B,0x2D,0x8B))

# Row 1 → Lambda collectors
box(s, Inches(2.55), Inches(1.5),  BW, BH, "health-collector", "Lambda Python 3.12", bg=NAVY)
box(s, Inches(2.55), Inches(3.0),  BW, BH, "deadline-reminder","Lambda Python 3.12", bg=NAVY)

# DynamoDB centre
box(s, Inches(5.0),  Inches(2.1),  BW, BH, "DynamoDB", "HealthEvents table", bg=RGBColor(0x1A,0x5F,0x3C))

# LLM summarizer
box(s, Inches(7.3),  Inches(1.5),  BW, BH, "llm-summarizer",   "Lambda Python 3.12", bg=NAVY)
box(s, Inches(9.5),  Inches(1.5),  BW, BH, "Amazon Bedrock",   "Claude Haiku",        bg=RGBColor(0x7B,0x2D,0x8B))

# SNS
box(s, Inches(5.0),  Inches(3.5),  BW, BH, "SNS Topic",        "Email digest",        bg=RGBColor(0xE8,0x55,0x1B))

# S3 archive
box(s, Inches(2.55), Inches(4.8),  BW, BH, "S3 Archive",       "Raw event JSON",      bg=RGBColor(0x56,0x7B,0x1C))

# Frontend stack
box(s, Inches(9.5),  Inches(3.5),  BW, BH, "React SPA",        "CloudFront + S3",     bg=RGBColor(0x00,0x52,0x9B))
box(s, Inches(9.5),  Inches(4.6),  BW, BH, "API Gateway",      "HTTP API",            bg=RGBColor(0xA0,0x52,0x2D))
box(s, Inches(7.3),  Inches(4.6),  BW, BH, "api-handler",      "Lambda Python 3.12",  bg=NAVY)

# Users
box(s, Inches(11.3), Inches(3.5),  Inches(1.6), BH, "👤 Users", "Browser", bg=RGBColor(0x44,0x44,0x44))

# arrows
cy = lambda top: top + Inches(0.36)   # vertical centre of a box
rx = lambda left: left + BW           # right edge
cx = lambda left: left + BW/2         # horizontal centre

arrow(s, rx(Inches(0.3)), cy(Inches(1.5)), Inches(2.55), cy(Inches(1.5)))
arrow(s, rx(Inches(0.3)), cy(Inches(3.0)), Inches(2.55), cy(Inches(3.0)))
arrow(s, rx(Inches(2.55)), cy(Inches(1.5)), Inches(5.0), cy(Inches(2.1)))
arrow(s, cx(Inches(5.0)), Inches(2.1), cx(Inches(5.0)), Inches(1.5)+Inches(0.3))  # up to summarizer row
arrow(s, Inches(5.0)+BW, cy(Inches(2.1)), Inches(7.3), cy(Inches(1.5)))
arrow(s, rx(Inches(7.3)), cy(Inches(1.5)), Inches(9.5), cy(Inches(1.5)))
arrow(s, rx(Inches(2.55)), cy(Inches(3.0)), Inches(5.0), cy(Inches(3.5)))
arrow(s, rx(Inches(2.55)), cy(Inches(1.5)), Inches(2.55)+BW/2, Inches(4.8))
arrow(s, rx(Inches(9.5)), cy(Inches(3.5)), Inches(11.3), cy(Inches(3.5)))
arrow(s, Inches(11.3)+Inches(0.8), cy(Inches(3.5)), Inches(9.5)+BW, cy(Inches(4.6)))
arrow(s, Inches(9.5), cy(Inches(4.6)), Inches(7.3)+BW, cy(Inches(4.6)))
arrow(s, Inches(7.3), cy(Inches(4.6)), Inches(5.0)+BW, cy(Inches(2.1))+Inches(0.1))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — AWS Services Table
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
navy_header(s, "AWS Services Used", "All FedRAMP Authorized  ·  us-east-1 commercial  ·  No NAT Gateways")
footer(s)

rows = [
    ("AWS Health API (Orgs)", "Source of all health events across 200 accounts", "✅ Authorized", "~$0"),
    ("Lambda (×4)",           "health-collector, llm-summarizer, deadline-reminder, api-handler", "✅ Authorized", "~$1–2/mo"),
    ("DynamoDB (on-demand)",  "Primary store — events, status, follow-up tracking", "✅ Authorized", "~$5–10/mo"),
    ("Amazon Bedrock",        "Claude Haiku LLM — summaries & action extraction", "✅ Authorized", "~$1–2/mo"),
    ("S3 (×2)",               "Raw event archive + React SPA static assets", "✅ Authorized", "<$1/mo"),
    ("CloudFront",            "CDN for React SPA with Origin Access Control", "✅ Authorized", "<$1/mo"),
    ("API Gateway (HTTP)",    "REST backend for the React app", "✅ Authorized", "~$1/mo"),
    ("EventBridge Scheduler", "Cron triggers: 15-min collector + daily reminder", "✅ Authorized", "<$1/mo"),
    ("SNS",                   "Deadline alert email digests to ops team", "✅ Authorized", "<$1/mo"),
    ("CloudWatch Logs",       "All Lambda logs, 30-day retention", "✅ Authorized", "<$1/mo"),
]

headers = ["Service", "Purpose", "FedRAMP", "Est. Cost"]
col_w   = [Inches(2.2), Inches(6.6), Inches(1.6), Inches(1.6)]
col_x   = [Inches(0.3), Inches(2.5), Inches(9.1), Inches(10.7)]

# header row
y = Inches(1.45)
for i, h in enumerate(headers):
    add_rect(s, col_x[i], y, col_w[i], Inches(0.38), fill=NAVY)
    add_textbox(s, col_x[i]+Inches(0.06), y+Inches(0.05), col_w[i]-Inches(0.1), Inches(0.3),
                h, size=11, bold=True, color=WHITE)

y += Inches(0.38)
for ri, row in enumerate(rows):
    bg = LGRAY if ri % 2 == 0 else WHITE
    row_h = Inches(0.46)
    for ci, cell in enumerate(row):
        add_rect(s, col_x[ci], y, col_w[ci], row_h, fill=bg, line=RGBColor(0xDD,0xDD,0xDD))
        clr = GREEN if cell.startswith("✅") else DGRAY
        add_textbox(s, col_x[ci]+Inches(0.06), y+Inches(0.07), col_w[ci]-Inches(0.1), row_h-Inches(0.1),
                    cell, size=10, color=clr)
    y += row_h

add_textbox(s, Inches(0.3), Inches(7.0), Inches(10), Inches(0.35),
            "Total estimated monthly cost: ~$10–15/month at steady state (50K events, 500 Bedrock calls)",
            size=11, bold=True, color=NAVY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Dashboard Screenshot
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
navy_header(s, "Dashboard", "Real-time stats, urgency charts, and upcoming deadlines at a glance")
footer(s)

img_path = os.path.join(SCREENS, "01_dashboard.png")
if os.path.exists(img_path):
    s.shapes.add_picture(img_path, Inches(0.3), Inches(1.35), Inches(12.7), Inches(5.65))

add_textbox(s, Inches(0.3), Inches(7.02), Inches(12.5), Inches(0.3),
            "Stats cards · Events by Service chart · Events by Urgency chart · Top critical events table",
            size=10, color=DGRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Notifications List Screenshot
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
navy_header(s, "Notifications List", "Searchable, filterable, sortable — all 200 accounts in one view")
footer(s)

img_path = os.path.join(SCREENS, "02_notification_list.png")
if os.path.exists(img_path):
    s.shapes.add_picture(img_path, Inches(0.3), Inches(1.35), Inches(12.7), Inches(5.65))

add_textbox(s, Inches(0.3), Inches(7.02), Inches(12.5), Inches(0.3),
            "Search · Filter by Service / Urgency / Status / Account · Deadline countdown · CSV export",
            size=10, color=DGRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Filtered View Screenshot
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
navy_header(s, "Intelligent Filtering", "Instantly isolate Critical events across all accounts")
footer(s)

img_path = os.path.join(SCREENS, "03_notification_list_filtered.png")
if os.path.exists(img_path):
    s.shapes.add_picture(img_path, Inches(0.3), Inches(1.35), Inches(12.7), Inches(5.65))

add_textbox(s, Inches(0.3), Inches(7.02), Inches(12.5), Inches(0.3),
            "Filter to Critical urgency — surface highest-priority items immediately for any account",
            size=10, color=DGRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Event Detail + AI Summary Screenshot
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
navy_header(s, "AI-Generated Event Detail", "Amazon Bedrock (Claude Haiku) turns raw API events into actionable guidance")
footer(s)

img_path = os.path.join(SCREENS, "04_notification_detail.png")
if os.path.exists(img_path):
    s.shapes.add_picture(img_path, Inches(0.3), Inches(1.35), Inches(8.2), Inches(5.65))

# call-out boxes on the right
callouts = [
    (BLUE,   "📝  Plain-English Summary",  "2–3 sentence stakeholder-ready description generated by Claude Haiku via Amazon Bedrock."),
    (GREEN,  "✅  Recommended Actions",     "3–5 specific, step-by-step action items extracted from the raw event description."),
    (RED,    "⏰  Urgency + Deadline",      "Automatic classification: Critical / High / Medium / Low based on deadline proximity."),
    (NAVY,   "👤  Follow-up Tracker",      "Owner assignment, status dropdown, and notes — saved directly to DynamoDB."),
]

y = Inches(1.45)
for bg, title, body in callouts:
    add_rect(s, Inches(8.7), y, Inches(4.35), Inches(1.28), fill=bg)
    add_textbox(s, Inches(8.85), y+Inches(0.08), Inches(4.1), Inches(0.38),
                title, size=12, bold=True, color=WHITE)
    add_textbox(s, Inches(8.85), y+Inches(0.45), Inches(4.1), Inches(0.75),
                body, size=10, color=WHITE, wrap=True)
    y += Inches(1.38)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Security & Compliance
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
navy_header(s, "Security & Compliance", "FedRAMP Authorized services only · IAM least-privilege · No data leaves us-east-1")
footer(s)

left_items = [
    ("🏛  FedRAMP Authorized",   "All AWS services used (Lambda, DynamoDB, S3, CloudFront, API Gateway, Bedrock, EventBridge, SNS) are FedRAMP Authorized in us-east-1."),
    ("🔐  IAM Least-Privilege",  "Each Lambda has a dedicated IAM role with only the permissions it needs — no wildcard actions, no shared roles."),
    ("🔒  Encryption at Rest",   "DynamoDB tables encrypted with AWS-managed keys. S3 buckets use S3-managed encryption. All traffic over TLS 1.2+."),
    ("🌐  Network Isolation",    "No VPC, no NAT Gateway — Lambda functions call AWS APIs directly over HTTPS. No public S3 bucket policies; CloudFront uses OAC."),
    ("📋  Audit & Observability","All Lambda invocations logged to CloudWatch with 30-day retention. DynamoDB Point-in-Time Recovery enabled on the HealthEvents table."),
]

y = Inches(1.45)
for title, body in left_items:
    add_rect(s, Inches(0.3), y, Inches(12.6), Inches(0.95), fill=LGRAY, line=RGBColor(0xDD,0xDD,0xDD))
    add_rect(s, Inches(0.3), y, Inches(0.12), Inches(0.95), fill=ORANGE)
    add_textbox(s, Inches(0.55), y+Inches(0.06), Inches(3.2), Inches(0.38),
                title, size=13, bold=True, color=NAVY)
    add_textbox(s, Inches(3.8), y+Inches(0.14), Inches(9.0), Inches(0.65),
                body, size=11, color=DGRAY, wrap=True)
    y += Inches(1.05)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Cost Summary
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
navy_header(s, "Cost Estimate", "Pay-per-use serverless architecture — minimal idle cost")
footer(s)

cost_rows = [
    ("Lambda (×4 functions)",  "15-min polling + stream triggers + daily reminder + API calls", "~$1–2/mo"),
    ("DynamoDB on-demand",     "~50K reads + 50K writes per month (200 accounts × daily events)", "~$5–10/mo"),
    ("Amazon Bedrock (Haiku)", "~500 events/month × avg 1K tokens in + 0.5K out", "~$1–2/mo"),
    ("S3 (archive + SPA)",     "Raw JSON archive + static frontend assets", "<$1/mo"),
    ("CloudFront",             "Static SPA delivery — free tier covers most traffic", "<$1/mo"),
    ("API Gateway HTTP API",   "~10K API calls/month from dashboard users", "~$1/mo"),
    ("EventBridge Scheduler",  "2 schedules — flat rate", "<$1/mo"),
    ("SNS + CloudWatch",       "Email digest + log ingestion", "<$1/mo"),
]

col_x2 = [Inches(0.3), Inches(2.9), Inches(10.3)]
col_w2 = [Inches(2.5), Inches(7.3), Inches(2.7)]
hdrs   = ["Service", "Basis", "Monthly Cost"]

y = Inches(1.45)
for i, h in enumerate(hdrs):
    add_rect(s, col_x2[i], y, col_w2[i], Inches(0.38), fill=NAVY)
    add_textbox(s, col_x2[i]+Inches(0.06), y+Inches(0.05), col_w2[i]-Inches(0.1), Inches(0.3),
                h, size=11, bold=True, color=WHITE)
y += Inches(0.38)
for ri, (svc, basis, cost) in enumerate(cost_rows):
    bg = LGRAY if ri % 2 == 0 else WHITE
    rh = Inches(0.44)
    for ci, cell in enumerate([svc, basis, cost]):
        add_rect(s, col_x2[ci], y, col_w2[ci], rh, fill=bg, line=RGBColor(0xDD,0xDD,0xDD))
        clr = GREEN if "mo" in cell and cell.strip().startswith("~") or cell.startswith("<") else DGRAY
        add_textbox(s, col_x2[ci]+Inches(0.06), y+Inches(0.07), col_w2[ci]-Inches(0.1), rh-Inches(0.1),
                    cell, size=10, color=clr)
    y += rh

# total bar
add_rect(s, Inches(0.3), y+Inches(0.08), Inches(12.7), Inches(0.5), fill=NAVY)
add_textbox(s, Inches(0.4), y+Inches(0.12), Inches(9), Inches(0.36),
            "Total Estimated Monthly Cost", size=13, bold=True, color=WHITE)
add_textbox(s, Inches(10.0), y+Inches(0.12), Inches(3), Inches(0.36),
            "~$10–15 / month", size=14, bold=True, color=ORANGE, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Next Steps & Roadmap
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
navy_header(s, "Roadmap & Next Steps", "Phase 2 enhancements to further reduce operational toil")
footer(s)

roadmap = [
    ("Phase 1  ✅  Complete",
     ["Automated health event collection across all 200 accounts",
      "AI-generated summaries and action items via Amazon Bedrock",
      "Centralized React dashboard with filtering, search, and CSV export",
      "Follow-up ownership tracking and deadline reminders"]),
    ("Phase 2  →  Near-term",
     ["Slack / Microsoft Teams integration for real-time event notifications",
      "ServiceNow ticketing — auto-create incident for Critical events",
      "Cognito authentication when AWS SSO federation is available",
      "Multi-region support (us-west-2 failover)"]),
    ("Phase 3  →  Strategic",
     ["AWS Systems Manager integration — automated remediation runbooks",
      "Account tag-based routing (team ownership from AWS Organizations tags)",
      "Historical trend analysis — Bedrock insights on recurring event patterns",
      "API integration with USPTO ITSM / CMDB for asset correlation"]),
]

y = Inches(1.42)
for col_idx, (phase, bullets) in enumerate(roadmap):
    x = Inches(0.3) + col_idx * Inches(4.35)
    bg = GREEN if "Complete" in phase else (NAVY if "Near" in phase else RGBColor(0x56,0x56,0x56))
    add_rect(s, x, y, Inches(4.15), Inches(0.45), fill=bg)
    add_textbox(s, x+Inches(0.1), y+Inches(0.06), Inches(3.95), Inches(0.35),
                phase, size=12, bold=True, color=WHITE)
    add_rect(s, x, y+Inches(0.45), Inches(4.15), Inches(5.35), fill=LGRAY, line=RGBColor(0xDD,0xDD,0xDD))
    by = y + Inches(0.6)
    for b in bullets:
        add_textbox(s, x+Inches(0.15), by, Inches(3.85), Inches(0.55),
                    f"• {b}", size=11, color=DGRAY, wrap=True)
        by += Inches(1.15)

add_textbox(s, Inches(0.3), Inches(7.02), Inches(12.5), Inches(0.3),
            "Questions?  |  Repository: github.com/saspto/aws-health-notifications-tracker  |  Live app: d2vuz8nfxuvkbe.cloudfront.net",
            size=10, color=DGRAY, align=PP_ALIGN.CENTER)


# ── save ────────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"Saved: {OUT}")
